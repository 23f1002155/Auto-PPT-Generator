import streamlit as st
import openai
from pptx import Presentation
import io
import json

# --- Core Functions ---

def generate_slide_content(api_key, raw_text, guidance):
    """
    Uses an LLM to structure raw text into a JSON format for slides.
    NOTE: This example uses OpenAI. You can adapt it for other providers.
    """
    openai.api_key = api_key
    
    prompt = f"""
    You are an expert at structuring long-form text into concise and clear presentation slides.
    Your task is to analyze the provided text and guidance to create the content for a slide deck.

    **Guidance:** "{guidance}"

    **Text to Analyze:**
    \"\"\"
    {raw_text}
    \"\"\"

    Based on the text and guidance, break it down into a logical sequence of slides.
    Determine a reasonable number of slides needed to cover the content effectively.

    **Output Format:**
    Return ONLY a valid JSON object with a single key "slides". Each item in the "slides" array
    should be an object representing one slide, containing a "title" and a "content" array of strings
    (for bullet points or short paragraphs).

    **Example JSON Output:**
    {{
      "slides": [
        {{
          "title": "Slide 1 Title",
          "content": [
            "First bullet point on slide 1.",
            "Second bullet point on slide 1."
          ]
        }},
        {{
          "title": "Slide 2 Title",
          "content": [
            "A single paragraph of text for this slide."
          ]
        }}
      ]
    }}
    """

    try:
        response = openai.chat.completions.create(
            model="gpt-3.5-turbo", # Or any other suitable model
            messages=[
                {"role": "system", "content": "You are a helpful presentation assistant."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.5,
        )
        # Extract the JSON string from the response
        json_string = response.choices[0].message.content
        return json.loads(json_string)
    except Exception as e:
        st.error(f"An error occurred with the LLM API call: {e}")
        return None

def create_presentation(llm_output_json, template_file):
    """
    Creates a PowerPoint presentation from the LLM output using a template.
    """
    try:
        prs = Presentation(template_file)
        
        # Use the "Title and Content" layout (usually index 1, but can vary)
        # A more robust solution would be to find the layout by its name
        title_and_content_layout = prs.slide_layouts[1]

        for slide_data in llm_output_json.get('slides', []):
            slide = prs.slides.add_slide(title_and_content_layout)
            
            title = slide.shapes.title
            if title:
                title.text = slide_data.get('title', '')
            
            # Find the main content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    content_placeholder = shape
                    break

            if content_placeholder:
                tf = content_placeholder.text_frame
                tf.clear()  # Clear existing default text
                
                content_items = slide_data.get('content', [])
                if content_items:
                    # Set the first item without a bullet
                    p = tf.paragraphs[0]
                    p.text = content_items[0]
                    # Add remaining items as bullet points
                    for item in content_items[1:]:
                        p = tf.add_paragraph()
                        p.text = item
                        p.level = 0
            
        # Save the presentation to an in-memory stream
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        return ppt_stream
    except Exception as e:
        st.error(f"Failed to create presentation file: {e}")
        return None


# --- Streamlit UI ---

st.set_page_config(layout="centered", page_title="Auto PPT Generator")
st.title("✨ Your Text, Your Style – Auto PPT Generator")

st.markdown("""
Create a fully formatted PowerPoint presentation from bulk text, markdown, or prose, perfectly matching your chosen template's look and feel.
""")

# 1. Text Input
with st.expander("1. Paste your text content", expanded=True):
    raw_text = st.text_area("Paste your content here", height=250, label_visibility="collapsed")

# 2. Optional Guidance
with st.expander("2. Provide guidance (Optional)"):
    guidance = st.text_input("Enter a one-line guidance for the tone or structure", placeholder="e.g., turn into an investor pitch deck")

# 3. API Key Input
with st.expander("3. Enter your LLM API Key", expanded=True):
    api_key = st.text_input("Provide your LLM API key (e.g. OpenAI)", type="password", label_visibility="collapsed")
    st.info("Your API key is used only for this session and is never stored or logged.", icon="🔒")

# 4. Template Upload
with st.expander("4. Upload your PowerPoint Template", expanded=True):
    uploaded_template = st.file_uploader("Upload a PowerPoint template or presentation (.pptx or .potx)", type=['pptx', 'potx'], label_visibility="collapsed")


# 5. Generate Button
st.divider()
generate_button = st.button("🚀 Generate Presentation", type="primary", use_container_width=True)

if generate_button:
    # --- Input Validation ---
    if not raw_text:
        st.error("Please paste your text content in Step 1.")
    elif not api_key:
        st.error("Please enter your API key in Step 3.")
    elif not uploaded_template:
        st.error("Please upload a PowerPoint template in Step 4.")
    else:
        with st.spinner("Generating slide content with AI... 🧠"):
            llm_output = generate_slide_content(api_key, raw_text, guidance)
        
        if llm_output:
            st.success("✅ Content structure generated!")
            with st.spinner("Creating your presentation... 🎨"):
                generated_ppt_stream = create_presentation(llm_output, uploaded_template)
            
            if generated_ppt_stream:
                st.success("✅ Presentation created!")
                
                st.download_button(
                    label="📥 Download Presentation",
                    data=generated_ppt_stream,
                    file_name="generated_presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )